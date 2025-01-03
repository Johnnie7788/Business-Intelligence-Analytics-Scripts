#!/usr/bin/env python
# coding: utf-8

# In[ ]:


import pandas as pd
import matplotlib.pyplot as plt
from sklearn.metrics.pairwise import cosine_similarity
import logging
import json
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class StrategicDataTransformation:
    def __init__(self, analytics_data_path, tools_comparison_path):
        """Initialize the class with paths to analytics infrastructure data and tools comparison data."""
        self.analytics_data = pd.read_csv(analytics_data_path)
        self.tools_comparison = pd.read_csv(tools_comparison_path)
        self.transformation_plan = {}

    def assess_analytics_infrastructure(self):
        """Assess the current analytics infrastructure and identify gaps."""
        logging.info("Assessing analytics infrastructure...")
        try:
            infrastructure_summary = self.analytics_data.describe(include='all').to_dict()
            gaps = {
                "missing_sources": "Identify data sources not integrated into pipelines.",
                "processing_delays": "Evaluate latency in data processing.",
                "scalability": "Assess the ability to handle increased data volume."
            }
            self.transformation_plan['current_infrastructure'] = infrastructure_summary
            self.transformation_plan['identified_gaps'] = gaps
            logging.info("Analytics infrastructure assessment completed.")
        except Exception as e:
            logging.error(f"Error during analytics infrastructure assessment: {e}")
            raise

    def propose_transformation_plan(self):
        """Propose a strategic data transformation roadmap."""
        logging.info("Proposing transformation plan...")
        try:
            roadmap = {
                "short_term": "Consolidate data sources and standardize data pipelines.",
                "medium_term": "Implement advanced analytics with predictive modeling capabilities.",
                "long_term": "Integrate AI-driven insights, real-time data processing, and cloud scalability."
            }
            self.transformation_plan['roadmap'] = roadmap
            logging.info("Transformation plan proposed successfully.")
        except Exception as e:
            logging.error(f"Error during transformation plan proposal: {e}")
            raise

    def evaluate_tools(self):
        """Compare and recommend BI tools based on organizational needs."""
        logging.info("Evaluating BI tools...")
        try:
            self.tools_comparison['average_score'] = self.tools_comparison.mean(axis=1)
            recommended_tools = self.tools_comparison.sort_values(by='average_score', ascending=False)
            self.transformation_plan['tools_evaluation'] = recommended_tools[['Tool', 'average_score']].to_dict('records')

            plt.figure(figsize=(10, 6))
            plt.bar(recommended_tools['Tool'], recommended_tools['average_score'], color='skyblue')
            plt.title("BI Tools Evaluation")
            plt.xlabel("Tools")
            plt.ylabel("Average Score")
            plt.xticks(rotation=45)
            plt.tight_layout()
            plt.show()

            logging.info("BI tools evaluation completed successfully.")
        except Exception as e:
            logging.error(f"Error during tools evaluation: {e}")
            raise

    def create_presentations(self, output_path):
        """Generate presentations and visualizations to communicate the strategy."""
        logging.info("Creating presentation and visualizations...")
        try:
            presentation_path = f"{output_path}/transformation_plan_{pd.Timestamp.now().strftime('%Y%m%d%H%M%S')}.json"
            with open(presentation_path, 'w') as file:
                json.dump(self.transformation_plan, file, indent=4)

            ppt_path = f"{output_path}/transformation_strategy_{pd.Timestamp.now().strftime('%Y%m%d%H%M%S')}.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Strategic Data Transformation Plan"

            for section, content in self.transformation_plan.items():
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = section.replace('_', ' ').title()
                text_box = slide.shapes.add_textbox(100, 100, 600, 400)
                frame = text_box.text_frame
                if isinstance(content, dict):
                    for key, value in content.items():
                        frame.add_paragraph().text = f"{key}: {value}"
                elif isinstance(content, list):
                    for item in content:
                        frame.add_paragraph().text = str(item)
                else:
                    frame.add_paragraph().text = str(content)

            prs.save(ppt_path)
            logging.info(f"Transformation strategy saved as presentation to {ppt_path}")

        except Exception as e:
            logging.error(f"Error creating presentations: {e}")
            raise

if __name__ == "__main__":
    analytics_data_path = "analytics_infrastructure.csv"
    tools_comparison_path = "bi_tools_comparison.csv"
    output_path = "output"

    sdt = StrategicDataTransformation(analytics_data_path, tools_comparison_path)

    try:
        sdt.assess_analytics_infrastructure()
        sdt.propose_transformation_plan()
        sdt.evaluate_tools()
        sdt.create_presentations(output_path)
    except Exception as e:
        logging.error(f"An error occurred: {e}")

